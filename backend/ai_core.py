# --- START OF FILE ai_core.py ---

# Notebook/backend/ai_core.py
import os
import logging
import fitz  # PyMuPDF
import re
import config
# Near the top of ai_core.py
from langchain_community.vectorstores import FAISS
from langchain_ollama import OllamaEmbeddings, ChatOllama
# Removed incorrect OllamaLLM import if it was there from previous attempts
from langchain.text_splitter import RecursiveCharacterTextSplitter # <<<--- ENSURE THIS IS PRESENT
from langchain.docstore.document import Document
from langchain.chains import LLMChain
from langchain.prompts import PromptTemplate # Import PromptTemplate if needed directly here
from config import (
    OLLAMA_BASE_URL, OLLAMA_MODEL, OLLAMA_EMBED_MODEL, FAISS_FOLDER,
    DEFAULT_PDFS_FOLDER, UPLOAD_FOLDER, RAG_CHUNK_K, MULTI_QUERY_COUNT,
    ANALYSIS_MAX_CONTEXT_LENGTH, OLLAMA_REQUEST_TIMEOUT, RAG_SEARCH_K_PER_QUERY,
    SUB_QUERY_PROMPT_TEMPLATE, SYNTHESIS_PROMPT_TEMPLATE, ANALYSIS_PROMPTS, db
)
from utils import parse_llm_response, escape_html # Added escape_html for potential use
import pdfplumber
import gridfs
import tempfile
import requests
from docx import Document
from pptx import Presentation

fs = gridfs.GridFS(db)

logger = logging.getLogger(__name__)

# --- Global State (managed within functions) ---
document_texts_cache = {}
vector_store = None
embeddings: OllamaEmbeddings | None = None
llm: ChatOllama | None = None

# --- Initialization Functions ---

# ai_core.py (only showing the modified function)
def initialize_ai_components() -> tuple[OllamaEmbeddings | None, ChatOllama | None]:
    """Initializes Ollama Embeddings and LLM instances globally.

    Returns:
        tuple[OllamaEmbeddings | None, ChatOllama | None]: The initialized embeddings and llm objects,
                                                          or (None, None) if initialization fails.
    """
    global embeddings, llm
    if embeddings and llm:
        logger.info("AI components already initialized.")
        return embeddings, llm

    try:
        # Use the new OllamaEmbeddings from langchain_ollama
        logger.info(f"Initializing Ollama Embeddings: model={OLLAMA_EMBED_MODEL}, base_url={OLLAMA_BASE_URL}, timeout={OLLAMA_REQUEST_TIMEOUT}s")
        embeddings = OllamaEmbeddings(
            model=OLLAMA_EMBED_MODEL,
            base_url=OLLAMA_BASE_URL,
            #request_timeout=OLLAMA_REQUEST_TIMEOUT # Explicitly pass timeout
        )
        # Perform a quick test embedding
        _ = embeddings.embed_query("Test embedding query")
        logger.info("Ollama Embeddings initialized successfully.")

        # Use the new ChatOllama from langchain_ollama
        logger.info(f"Initializing Ollama LLM: model={OLLAMA_MODEL}, base_url={OLLAMA_BASE_URL}, timeout={OLLAMA_REQUEST_TIMEOUT}s")
        llm = ChatOllama(
            model=OLLAMA_MODEL,
            base_url=OLLAMA_BASE_URL,
            #request_timeout=OLLAMA_REQUEST_TIMEOUT # Explicitly pass timeout
        )
        # Perform a quick test invocation
        _ = llm.invoke("Respond briefly with 'AI Check OK'")
        logger.info("Ollama LLM initialized successfully.")

        return embeddings, llm  # Return the objects
    except ImportError as e:
        logger.critical(f"Import error during AI initialization: {e}. Ensure correct langchain packages are installed.", exc_info=True)
        embeddings = None
        llm = None
        return None, None
    except Exception as e:
        # Catch potential Pydantic validation error specifically if possible, or general Exception
        logger.error(f"Failed to initialize AI components (check Ollama server status, model name '{OLLAMA_MODEL}' / '{OLLAMA_EMBED_MODEL}', base URL '{OLLAMA_BASE_URL}', timeout {OLLAMA_REQUEST_TIMEOUT}s): {e}", exc_info=True)
        # Log the type of error for better debugging
        logger.error(f"Error Type: {type(e).__name__}")
        # If it's a Pydantic error, the message usually contains details
        if "pydantic" in str(type(e)).lower():
             logger.error(f"Pydantic Validation Error Details: {e}")
        embeddings = None
        llm = None
        return None, None

def load_vector_store() -> bool:
    """Loads the FAISS index from disk into the global `vector_store`.

    Requires `embeddings` to be initialized first.

    Returns:
        bool: True if the index was loaded successfully, False otherwise (or if not found).
    """
    global vector_store, embeddings
    if vector_store:
        logger.info("Vector store already loaded.")
        return True
    if not embeddings:
        logger.error("Embeddings not initialized. Cannot load vector store.")
        return False

    faiss_index_path = os.path.join(FAISS_FOLDER, "index.faiss")
    faiss_pkl_path = os.path.join(FAISS_FOLDER, "index.pkl")

    if os.path.exists(faiss_index_path) and os.path.exists(faiss_pkl_path):
        try:
            logger.info(f"Loading FAISS index from folder: {FAISS_FOLDER}")
            # Note: Loading requires the same embedding model used for saving.
            # allow_dangerous_deserialization is required for FAISS/pickle
            vector_store = FAISS.load_local(
                folder_path=FAISS_FOLDER,
                embeddings=embeddings, # Pass the initialized embeddings object
                allow_dangerous_deserialization=True
            )
            index_size = getattr(getattr(vector_store, 'index', None), 'ntotal', 0)
            if index_size > 0:
                logger.info(f"FAISS index loaded successfully. Contains {index_size} vectors.")
                return True
            else:
                logger.warning(f"FAISS index loaded from {FAISS_FOLDER}, but it appears to be empty.")
                return True # Treat empty as loaded
        except FileNotFoundError:
            logger.warning(f"FAISS index files not found in {FAISS_FOLDER}, although directory exists. Proceeding without loaded index.")
            vector_store = None
            return False
        except EOFError:
            logger.error(f"EOFError loading FAISS index from {FAISS_FOLDER}. Index file might be corrupted or incomplete.", exc_info=True)
            vector_store = None
            return False
        except Exception as e:
            logger.error(f"Error loading FAISS index from {FAISS_FOLDER}: {e}", exc_info=True)
            vector_store = None # Ensure it's None if loading failed
            return False
    else:
        logger.warning(f"FAISS index files (index.faiss, index.pkl) not found at {FAISS_FOLDER}. Will be created on first upload or if default.py ran.")
        vector_store = None
        return False # Indicate index wasn't loaded


def save_vector_store() -> bool:
    """Saves the current global `vector_store` (FAISS index) to disk.

    Returns:
        bool: True if saving was successful, False otherwise (or if store is None).
    """
    global vector_store
    if not vector_store:
        logger.warning("Attempted to save vector store, but it's not loaded or initialized.")
        return False
    if not os.path.exists(FAISS_FOLDER):
        try:
            os.makedirs(FAISS_FOLDER)
            logger.info(f"Created FAISS store directory: {FAISS_FOLDER}")
        except OSError as e:
            logger.error(f"Failed to create FAISS store directory {FAISS_FOLDER}: {e}", exc_info=True)
            return False

    try:
        index_size = getattr(getattr(vector_store, 'index', None), 'ntotal', 0)
        logger.info(f"Saving FAISS index ({index_size} vectors) to {FAISS_FOLDER}...")
        vector_store.save_local(FAISS_FOLDER)
        logger.info(f"FAISS index saved successfully.")
        return True
    except Exception as e:
        logger.error(f"Error saving FAISS index to {FAISS_FOLDER}: {e}", exc_info=True)
        return False


def load_all_document_texts():
    """Loads text from all supported files, embeds it, and adds it to the FAISS vector store."""
    global document_texts_cache
    document_texts_cache = {}  # Reset cache
    supported_extensions = ['.pdf', '.docx', '.pptx', '.txt']

    def _load_from_folder(folder_path):
        if not os.path.exists(folder_path):
            logger.warning(f"Folder not found: {folder_path}")
            return
        for filename in os.listdir(folder_path):
            ext = os.path.splitext(filename)[1].lower()
            file_path = os.path.join(folder_path, filename)
            try:
                if ext == '.pdf':
                    markdown = convert_pdf_to_markdown(file_path)
                elif ext == '.docx':
                    markdown = convert_docx_to_markdown(file_path)
                elif ext == '.pptx':
                    markdown = convert_pptx_to_markdown(file_path)
                elif ext == '.txt':
                    markdown = convert_txt_to_markdown(file_path)
                else:
                    continue

                # Embed and add Markdown data to FAISS
                if markdown:
                    success = add_markdown_to_vector_store(markdown, filename)
                    if not success:
                        logger.error(f"Failed to add embedded Markdown data to vector store for file: {filename}")
                document_texts_cache[filename] = markdown
            except Exception as e:
                logger.error(f"Error processing file {filename}: {e}", exc_info=True)

    _load_from_folder(config.DEFAULT_PDFS_FOLDER)
    _load_from_folder(config.UPLOAD_FOLDER)
    logger.info(f"Loaded {len(document_texts_cache)} documents into cache.")

# --- PDF Processing Functions ---

def extract_text_from_pdf(pdf_path: str) -> str | None:
    """Extracts text from a single PDF file and converts it to Markdown."""
    text = ""
    if not os.path.exists(pdf_path):
        logger.error(f"PDF file not found for extraction: {pdf_path}")
        return None
    try:
        with pdfplumber.open(pdf_path) as pdf:
            for page in pdf.pages:
                page_text = page.extract_text() or ""
                text += page_text + "\n\n"
        return convert_to_markdown(text.strip(), os.path.basename(pdf_path))
    except Exception as e:
        logger.error(f"Error extracting text from PDF {pdf_path}: {e}", exc_info=True)
        return None

# def extract_text_from_docx(docx_path: str) -> str | None:
#     """Extracts text from a DOCX file and converts it to Markdown."""
#     try:
#         doc = Document(docx_path)
#         text = "\n\n".join([p.text for p in doc.paragraphs if p.text.strip()])
#         return convert_to_markdown(text, os.path.basename(docx_path))
#     except Exception as e:
#         logger.error(f"Error extracting text from DOCX {docx_path}: {e}", exc_info=True)
#         return None

# def extract_text_from_pptx(pptx_path: str) -> str | None:
#     """Extracts text from a PPTX file and converts it to Markdown."""
#     try:
#         presentation = Presentation(pptx_path)
#         text = ""
#         for slide in presentation.slides:
#             for shape in slide.shapes:
#                 if shape.has_text_frame:
#                     text += shape.text_frame.text + "\n\n"
#         return convert_to_markdown(text, os.path.basename(pptx_path))
#     except Exception as e:
#         logger.error(f"Error extracting text from PPTX {pptx_path}: {e}", exc_info=True)
#         return None

# def extract_text_from_txt(txt_path: str) -> str | None:
#     """Extracts text from a TXT file and converts it to Markdown."""
#     try:
#         with open(txt_path, 'r', encoding='utf-8') as f:
#             text = f.read()
#         return convert_to_markdown(text, os.path.basename(txt_path))
#     except Exception as e:
#         logger.error(f"Error extracting text from TXT {txt_path}: {e}", exc_info=True)
#         return None

# def create_chunks_from_text(text: str, filename: str) -> list[Document]:
#     """Splits text into chunks using RecursiveCharacterTextSplitter and creates LangChain Documents.

#     Args:
#         text (str): The text content to chunk.
#         filename (str): The source filename for metadata.

#     Returns:
#         list[Document]: A list of LangChain Document objects representing the chunks.
#     """
#     if not text:
#         logger.warning(f"Cannot create chunks for '{filename}', input text is empty.")
#         return []

#     text_splitter = RecursiveCharacterTextSplitter(
#         chunk_size=1000,      # Target size of each chunk
#         chunk_overlap=150,    # Overlap between chunks
#         length_function=len,
#         add_start_index=True, # Include start index in metadata
#         separators=["\n\n", "\n", ". ", ", ", " ", ""], # Hierarchical separators
#     )

#     try:
#         # Use create_documents which handles metadata assignment more cleanly
#         documents = text_splitter.create_documents([text], metadatas=[{"source": filename}])
#         # Add explicit chunk_index for clarity (though start_index is also present)
#         for i, doc in enumerate(documents):
#             doc.metadata["chunk_index"] = i

#         logger.info(f"Created {len(documents)} LangChain Document chunks for '{filename}'.")
#         return documents

#     except Exception as e:
#         logger.error(f"Error creating chunks for '{filename}': {e}", exc_info=True)
#         return []

# def add_documents_to_vector_store(documents: list[Document]) -> bool:
#     """Adds Markdown documents to the vector store."""
#     global vector_store, embeddings
#     if not embeddings:
#         logger.error("Embeddings not initialized. Cannot add documents to vector store.")
#         return False

#     try:
#         if vector_store:
#             vector_store.add_documents(documents)
#         else:
#             vector_store = FAISS.from_documents(documents, embeddings)
#         return save_vector_store()
#     except Exception as e:
#         logger.error(f"Error adding documents to vector store: {e}", exc_info=True)
#         return False

from langchain.docstore.document import Document

def add_markdown_to_vector_store(markdown_text: str, filename: str) -> bool:
    """
    Adds Markdown data to the FAISS vector store.
    Args:
        markdown_text (str): The Markdown-formatted text.
        filename (str): The source filename for metadata.
    Returns:
        bool: True if successful, False otherwise.
    """
    global vector_store, embeddings
    if not embeddings:
        logger.error("Embeddings not initialized. Cannot add documents to vector store.")
        return False

    try:
        # Convert Markdown text into LangChain Document
        document = Document(page_content=markdown_text, metadata={"source": filename})
        if vector_store:
            vector_store.add_documents([document])  # Add document to existing vector store
        else:
            vector_store = FAISS.from_documents([document], embeddings)  # Create new vector store
        return save_vector_store()  # Save the updated vector store
    except Exception as e:
        logger.error(f"Error adding Markdown data to vector store: {e}", exc_info=True)
        return False# --- RAG and LLM Interaction ---

# --- MODIFIED: Added logging ---
def generate_sub_queries(query: str) -> list[str]:
    """
    Uses the LLM to generate sub-queries for RAG. Includes the original query.
    Uses SUB_QUERY_PROMPT_TEMPLATE from config.
    """
    global llm
    if not llm:
        logger.error("LLM not initialized, cannot generate sub-queries. Using original query only.")
        return [query]
    if MULTI_QUERY_COUNT <= 0:
        logger.debug("MULTI_QUERY_COUNT is <= 0, skipping sub-query generation.")
        return [query]

    # Use the prompt template from config
    chain = LLMChain(llm=llm, prompt=SUB_QUERY_PROMPT_TEMPLATE)

    try:
        logger.info(f"Generating {MULTI_QUERY_COUNT} sub-queries for: '{query[:100]}...'")
        # Log the prompt before sending (approx first 150 chars)
        prompt_to_log = SUB_QUERY_PROMPT_TEMPLATE.format(query=query, num_queries=MULTI_QUERY_COUNT)
        logger.debug(f"Sub-query Prompt (Start):\n{prompt_to_log[:150]}...") # DEBUG level might be better

        response = chain.invoke({"query": query, "num_queries": MULTI_QUERY_COUNT})
        # Response structure might vary; often {'text': 'query1\nquery2'}
        raw_response_text = response.get('text', '') if isinstance(response, dict) else str(response)

        # Log the raw response start
        logger.debug(f"Sub-query Raw Response (Start):\n{raw_response_text[:150]}...") # DEBUG level

        # No need to parse for <thinking> here as the prompt doesn't request it
        sub_queries = [q.strip() for q in raw_response_text.strip().split('\n') if q.strip()]

        if sub_queries:
            logger.info(f"Generated {len(sub_queries)} sub-queries.")
            # Ensure we don't exceed MULTI_QUERY_COUNT, and always include the original
            final_queries = [query] + sub_queries[:MULTI_QUERY_COUNT]
            # Deduplicate the final list just in case LLM generated the original query
            final_queries = list(dict.fromkeys(final_queries))
            logger.debug(f"Final search queries: {final_queries}")
            return final_queries
        else:
            logger.warning("LLM did not generate any valid sub-queries. Falling back to original query only.")
            return [query]

    except Exception as e:
        logger.error(f"Error generating sub-queries: {e}", exc_info=True)
        return [query] # Fallback
# --- END MODIFICATION ---

def perform_rag_search(query: str, user_email=None) -> tuple[list[Document], str, dict[int, dict]]:
    """
    Performs RAG using Markdown data stored in the FAISS vector store.
    Args:
        query (str): The user query.
    Returns:
        tuple[list[Document], str, dict[int, dict]]: 
            - List of retrieved Document objects.
            - Formatted context text for the LLM prompt.
            - Citation map for frontend reference.
    """
    global vector_store
    if not vector_store:
        logger.warning("RAG search attempted but no vector store is loaded.")
        return [], "No context available.", {}

    if not query or not query.strip():
        logger.warning("RAG search attempted with empty query.")
        return [], "No context available.", {}

    index_size = getattr(getattr(vector_store, 'index', None), 'ntotal', 0)
    if index_size == 0:
        logger.warning("RAG search attempted but the vector store index is empty.")
        return [], "No context available.", {}

    try:
        # Step 1: Generate Sub-Queries
        search_queries = generate_sub_queries(query)
        logger.info(f"Generated {len(search_queries)} sub-queries for query: '{query[:100]}...'")

        # Step 2: Perform Similarity Search for each sub-query
        all_retrieved_docs_with_scores = []
        for sub_query in search_queries:
            results_with_scores = vector_store.similarity_search_with_score(sub_query, k=config.RAG_SEARCH_K_PER_QUERY)
            all_retrieved_docs_with_scores.extend(results_with_scores)

        # Step 3: Deduplicate and Sort Documents by Score
        unique_docs_dict = {}
        for doc, score in all_retrieved_docs_with_scores:
            source = doc.metadata.get('source', 'Unknown Source')
            content_hash = hash(doc.page_content)
            doc_key = (source, content_hash)

            # Keep the document with the best score (lowest distance)
            if doc_key not in unique_docs_dict or score < unique_docs_dict[doc_key][1]:
                unique_docs_dict[doc_key] = (doc, score)

        # Sort unique documents by score (ascending - best first)
        sorted_unique_docs = sorted(unique_docs_dict.values(), key=lambda item: item[1])

        # Select the final top RAG_CHUNK_K unique documents
        final_context_docs_with_scores = sorted_unique_docs[:config.RAG_CHUNK_K]
        context_docs = [doc for doc, score in final_context_docs_with_scores]

        # logger.info(f"Retrieved {len(all_retrieved_docs_with_scores)} chunks total across sub-queries. "
        #             f"Selected {len(context_docs)} unique chunks (target k={config.RAG_CHUNK_K}) for context.")

        # Step 4: Format Context for LLM Prompt and Create Citation Map
        formatted_context_parts = []
        context_docs_map = {}  # Use 1-based index for map keys, matching citations like [1], [2]
        for i, doc in enumerate(context_docs):
            citation_index = i + 1  # 1-based index for the prompt and map
            source = doc.metadata.get('source', 'Unknown Source')
            chunk_idx = doc.metadata.get('chunk_index', 'N/A')
            content = doc.page_content

            # Format for the LLM prompt
            context_str = f"[{citation_index}] Source: {source} | Chunk Index: {chunk_idx}\n{content}"
            formatted_context_parts.append(context_str)

            # Store data needed for frontend reference display, keyed by the citation number
            context_docs_map[citation_index] = {
                "source": source,
                "chunk_index": chunk_idx,
                "content": content
            }

        formatted_context_text = "\n\n---\n\n".join(formatted_context_parts) if formatted_context_parts else "No context chunks selected after processing."

    except Exception as e:
        logger.error(f"Error during RAG search process for query '{query[:50]}...': {e}", exc_info=True)
        # Reset results on error
        context_docs = []
        formatted_context_text = "Error retrieving context due to an internal server error."
        context_docs_map = {}

    # Return the list of Document objects, the formatted text for the LLM, and the citation map
    return context_docs, formatted_context_text, context_docs_map
# --- MODIFIED: Added logging ---
def synthesize_chat_response(query: str, context_text: str) -> tuple[str, str | None]:
    """
    Generates the final chat response using the LLM, query, and context.
    Requests and parses thinking/reasoning content using SYNTHESIS_PROMPT_TEMPLATE.

    Returns:
        tuple[str, str | None]: (user_answer, thinking_content)
    """
    global llm
    if not llm:
        logger.error("LLM not initialized, cannot synthesize response.")
        return "Error: The AI model is currently unavailable.", None

    # Use the prompt template from config
    # Ensure the prompt template is correctly formatted and expects 'query' and 'context'
    try:
        final_prompt = SYNTHESIS_PROMPT_TEMPLATE.format(query=query, context=context_text)
        # Log the prompt before sending
        logger.info(f"Sending synthesis prompt to LLM (model: {OLLAMA_MODEL})...")
        logger.debug(f"Synthesis Prompt (Start):\n{final_prompt[:200]}...") # Log more chars if needed

    except KeyError as e:
        logger.error(f"Error formatting SYNTHESIS_PROMPT_TEMPLATE: Missing key {e}. Check config.py.")
        return "Error: Internal prompt configuration issue.", None
    except Exception as e:
         logger.error(f"Error creating synthesis prompt: {e}", exc_info=True)
         return "Error: Could not prepare the request for the AI model.", None

    try:
        # logger.info(f"Invoking LLM for chat synthesis (model: {OLLAMA_MODEL})") # Already logged above
        # Use .invoke() for ChatOllama which returns AIMessage, access content with .content
        response_object = llm.invoke(final_prompt)
        # Ensure response_object has 'content' attribute
        full_llm_response = getattr(response_object, 'content', str(response_object))

        # Log the raw response start
        logger.info(f"LLM synthesis response received (length: {len(full_llm_response)}).")
        logger.debug(f"Synthesis Raw Response (Start):\n{full_llm_response[:200]}...")

        # Parse the response to separate thinking and answer using the utility function
        user_answer, thinking_content = parse_llm_response(full_llm_response)

        if thinking_content:
            logger.info(f"Parsed thinking content (length: {len(thinking_content)}).")
        else:
            # This is expected if the LLM didn't include the tags or the prompt was adjusted
            logger.debug("No <thinking> content found or parsed in the LLM response.")


        if not user_answer and thinking_content:
             logger.warning("Parsed user answer is empty after removing thinking block. The response might have only contained thinking.")
             # Decide how to handle this - return thinking as answer, or a specific message?
             # Let's return a message indicating this.
             user_answer = "[AI response consisted only of reasoning. No final answer provided. See thinking process.]"
        elif not user_answer and not thinking_content:
             logger.error("LLM response parsing resulted in empty answer and no thinking content.")
             user_answer = "[AI Response Processing Error: Empty result after parsing]"


        # Basic check if the answer looks like an error message generated by the LLM itself
        if user_answer.strip().startswith("Error:") or "sorry, I encountered an error" in user_answer.lower():
            logger.warning(f"LLM synthesis seems to have resulted in an error message: '{user_answer[:100]}...'")

        return user_answer.strip(), thinking_content # Return stripped answer and thinking

    except Exception as e:
        logger.error(f"LLM chat synthesis failed: {e}", exc_info=True)
        error_message = f"Sorry, I encountered an error while generating the response ({type(e).__name__}). The AI model might be unavailable, timed out, or failed internally."
        # Attempt to parse thinking even from error if possible? Unlikely to be useful.
        return error_message, None
# --- END MODIFICATION ---

# --- MODIFIED: Added logging ---
def generate_document_analysis(filename: str, analysis_type: str, user_email=None) -> tuple[str | None, str | None]:
    """
    Generates analysis (FAQ, Topics, Mindmap) for a specific document, optionally including thinking.
    Uses ANALYSIS_PROMPTS from config. Retrieves text from cache or disk.

    Returns:
        tuple[str | None, str | None]: (analysis_content, thinking_content)
                                    Returns (error_message, thinking_content) on failure.
                                    Returns (None, None) if document text cannot be found/loaded.
    """
    global llm, document_texts_cache
    logger.info(f"Starting analysis: type='{analysis_type}', file='{filename}'")

    if not llm:
        logger.error("LLM not initialized, cannot perform analysis.")
        return " Error: AI MODEL is not available for analysis.", None

    # --- Step 1: Get Document Text ---
    try:
        doc_text = get_document_text(filename, user_email=user_email)
    except FileNotFoundError as e:
        logger.error(str(e))
        return f"Error: Document '{filename}' not found.", None
    except Exception as e:
        logger.error(f"Unexpected error loading document '{filename}': {e}", exc_info=True)
        return f"Error: Failed to retrieve text content for '{filename}'.", None

    if not doc_text:
        logger.error(f"Analysis failed: doc_text is unexpectedly empty for '{filename}' after all checks.")
        return f"Error: Failed to retrieve text content for '{filename}'.", None


    # --- Step 2: Prepare Text for LLM (Truncation) ---
    original_length = len(doc_text)
    if original_length > ANALYSIS_MAX_CONTEXT_LENGTH:
        logger.warning(f"Document '{filename}' text too long ({original_length} chars), truncating to {ANALYSIS_MAX_CONTEXT_LENGTH} for '{analysis_type}' analysis.")
        # Truncate from the end, keeping the beginning
        doc_text_for_llm = doc_text[:ANALYSIS_MAX_CONTEXT_LENGTH]
        # Add a clear truncation marker
        doc_text_for_llm += "\n\n... [CONTENT TRUNCATED DUE TO LENGTH LIMIT]"
    else:
        doc_text_for_llm = doc_text
        logger.debug(f"Using full document text ({original_length} chars) for analysis '{analysis_type}'.")

    # --- Step 3: Get Analysis Prompt ---
    prompt_template = ANALYSIS_PROMPTS.get(analysis_type)
    if not prompt_template or not isinstance(prompt_template, PromptTemplate):
        logger.error(f"Invalid or missing analysis prompt template for type: {analysis_type} in config.py")
        return f"Error: Invalid analysis type '{analysis_type}' or misconfigured prompt.", None

    try:
        # Ensure the template expects 'doc_text_for_llm'
        final_prompt = prompt_template.format(doc_text_for_llm=doc_text_for_llm)
        # Log the prompt before sending
        logger.info(f"Sending analysis prompt to LLM (type: {analysis_type}, file: {filename}, model: {OLLAMA_MODEL})...")
        logger.debug(f"Analysis Prompt (Start):\n{final_prompt[:200]}...")

    except KeyError as e:
        logger.error(f"Error formatting ANALYSIS_PROMPTS[{analysis_type}]: Missing key {e}. Check config.py.")
        return f"Error: Internal prompt configuration issue for {analysis_type}.", None
    except Exception as e:
        logger.error(f"Error creating analysis prompt for {analysis_type}: {e}", exc_info=True)
        return f"Error: Could not prepare the request for the {analysis_type} analysis.", None


    # --- Step 4: Call LLM and Parse Response ---
    try:
        # logger.info(f"Invoking LLM for '{analysis_type}' analysis of '{filename}' (model: {OLLAMA_MODEL})") # Already logged above
        # Use .invoke() for ChatOllama
        response_object = llm.invoke(final_prompt)
        full_analysis_response = getattr(response_object, 'content', str(response_object))

        # Log the raw response start
        logger.info(f"LLM analysis response received for '{filename}' ({analysis_type}). Length: {len(full_analysis_response)}")
        logger.debug(f"Analysis Raw Response (Start):\n{full_analysis_response[:200]}...")

        # Parse potential thinking and main content using the utility function
        analysis_content, thinking_content = parse_llm_response(full_analysis_response)

        # Clean Mermaid code for mindmap analysis

        #######
        if analysis_type == "mindmap" and analysis_content:
            analysis_content = clean_mermaid_mindmap(analysis_content)
        elif analysis_type == "flowchart" and analysis_content:
            analysis_content = clean_mermaid_flowchart(analysis_content)

        if thinking_content:
            logger.info(f"Parsed thinking content from analysis response for '{filename}'.")
        # else: logger.debug(f"No thinking content found in analysis response for '{filename}'.") # Normal if not requested/provided

        if not analysis_content and thinking_content:
            logger.warning(f"Parsed analysis content is empty for '{filename}' ({analysis_type}). Response only contained thinking.")
            analysis_content = "[Analysis consisted only of reasoning. No final output provided. See thinking process.]"
        elif not analysis_content and not thinking_content:
            logger.error(f"LLM analysis response parsing resulted in empty content and no thinking for '{filename}' ({analysis_type}).")
            analysis_content = "[Analysis generation resulted in empty content after parsing.]"


        # Optional: Basic format validation could be added here if needed (e.g., check for Q:/A: in FAQ)

        logger.info(f"Analysis successful for '{filename}' ({analysis_type}).")
        return analysis_content.strip(), thinking_content # Return success tuple

    except Exception as e:
        logger.error(f"LLM analysis invocation error for {filename} ({analysis_type}): {e}", exc_info=True)
        # Try to return error message with thinking if parsing happened before error? Unlikely.
        return f"Error generating analysis: AI model failed ({type(e).__name__}). Check logs for details.", None
# --- END MODIFICATION ---

def clean_mermaid_code(llm_output: str) -> str:
    """
    Extracts and cleans Mermaid mindmap code from LLM output.
    Removes code block markers and extra text.
    """
    code = llm_output.strip()
    # Remove triple backticks and language hints
    if code.startswith("```"):
        code = code.lstrip("`")
        # Remove the first line (which may be 'mermaid')
        code = "\n".join(code.splitlines()[1:])
    if code.endswith("```"):
        code = code.rstrip("`")
    # Extract only the mindmap block if there is extra text
    if "mindmap" in code:
        code = code[code.index("mindmap"):]
    # Remove or truncate lines that are too long or have invalid characters
    cleaned_lines = []
    for line in code.splitlines():
        # Truncate node labels to 40 chars, remove problematic chars
        if "[" in line and "]" in line:
            parts = line.split("[")
            label = "[".join(parts[1:]).split("]")[0]
            label = label.replace("-", " ").replace("(", "").replace(")", "")
            if len(label) > 40:
                label = label[:37] + "..."
            line = parts[0] + "[" + label + "]"
        cleaned_lines.append(line)
    return "\n".join(cleaned_lines).strip()

def clean_mermaid_flowchart(llm_output: str) -> str:
    """Clean and sanitize Mermaid flowchart (graph TD) output from LLM."""
    lines = llm_output.strip().splitlines()

    # Remove triple backticks and whitespace
    lines = [line.strip("` ").rstrip() for line in lines if not line.strip().startswith("```")]

    # Find the line with "graph TD"
    header_index = next((i for i, line in enumerate(lines) if line.startswith("graph TD")), None)
    if header_index is None:
        return ""

    # Keep only lines from "graph TD" onward
    final_lines = lines[header_index:]

    # Deduplicate while preserving order
    seen = set()
    cleaned_lines = []
    for line in final_lines:
        if line and line not in seen:
            cleaned_lines.append(line)
            seen.add(line)

    return "\n".join(cleaned_lines)

def clean_mermaid_mindmap(llm_output: str) -> str:
    """
    Cleans Mermaid mindmap code from LLM output.
    Ensures 'mindmap' is on its own line and removes code block markers.
    """
    code = llm_output.strip()
    if code.startswith("```"):
        code = code.lstrip("`")
        code = "\n".join(code.splitlines()[1:])
    if code.endswith("```"):
        code = code.rstrip("`")
    idx = code.find('mindmap')
    if idx != -1:
        after = code[idx + len('mindmap'):]
        after = after.lstrip('\r\n\t :')
        if after and after[0] != '\n':
            after = '\n' + after
        code = 'mindmap' + after
    code = "\n".join([line for line in code.splitlines() if line.strip() != ""])
    return code.strip()

from config import db
import gridfs
import tempfile

fs = gridfs.GridFS(db)

def get_document_text(filename, user_email=None):
    # 1. Check cache first
    if filename in document_texts_cache:
        return document_texts_cache[filename]

    # 2. Try to load from default folder
    default_path = os.path.join(config.DEFAULT_PDFS_FOLDER, filename)
    if os.path.exists(default_path):
        text = extract_text_from_pdf(default_path)
        document_texts_cache[filename] = text
        return text

    # 3. Try to load from GridFS for user-uploaded files
    if user_email:
        file_obj = fs.find_one({"filename": filename, "metadata.user_email": user_email})
        if file_obj:
            with tempfile.NamedTemporaryFile(suffix=".pdf", delete=False) as temp:
                temp.write(file_obj.read())
                temp.flush()
                temp.close()  # Explicitly close the file to release locks
                try:
                    text = extract_text_from_pdf(temp.name)
                    document_texts_cache[filename] = text
                    if text:
                        logger.info(f"Successfully extracted text from GridFS PDF '{filename}' for user '{user_email}'.")
                    else:
                        logger.error(f"Failed to extract text from GridFS PDF '{filename}' for user '{user_email}'.")
                    return text
                finally:
                    os.remove(temp.name)  # Clean up the temporary file

    # 4. Not found
    raise FileNotFoundError(f"Document file '{filename}' not found in default folder or GridFS for user.")

GEMINI_API_KEY = os.getenv("GEMINI_API_KEY")
GEMINI_API_URL = os.getenv("GEMINI_API_URL", "https://generativelanguage.googleapis.com/v1/models/gemini-pro:generateContent")

def synthesize_gemini_response(query):
    """
    Sends the user query to Gemini API and returns the response.
    """
    if not GEMINI_API_KEY:
        return "Gemini API key not configured.", None

    headers = {
        "Content-Type": "application/json"
    }
    params = {
        "key": GEMINI_API_KEY
    }
    data = {
        "contents": [
            {
                "parts": [
                    {"text": query}
                ]
            }
        ]
    }
    try:
        response = requests.post(GEMINI_API_URL, headers=headers, params=params, json=data, timeout=30)
        response.raise_for_status()
        result = response.json()
        # Extract the main answer
        answer = result.get("candidates", [{}])[0].get("content", {}).get("parts", [{}])[0].get("text", "")
        return answer.strip() if answer else "Gemini did not return a response.", None
    except Exception as e:
        return f"Error contacting Gemini API: {e}", None

def convert_to_markdown(text: str, filename: str) -> str:
    """
    Converts plain text into Markdown format for RAG.
    Args:
        text (str): The extracted text content.
        filename (str): The source filename for metadata.
    Returns:
        str: Markdown-formatted text.
    """
    markdown = f"# {filename}\n\n{text.strip()}"
    return markdown

def convert_pdf_to_markdown(pdf_path: str) -> str:
    """Extracts text from a PDF file and converts it to Markdown."""
    text = ""
    if not os.path.exists(pdf_path):
        raise FileNotFoundError(f"PDF file not found: {pdf_path}")
    try:
        with pdfplumber.open(pdf_path) as pdf:
            for page in pdf.pages:
                page_text = page.extract_text() or ""
                text += page_text + "\n\n"
        markdown = f"# {os.path.basename(pdf_path)}\n\n{text.strip()}"
        return markdown
    except Exception as e:
        raise RuntimeError(f"Error extracting text from PDF {pdf_path}: {e}")

def convert_docx_to_markdown(docx_path: str) -> str:
    """Extracts text from a DOCX file and converts it to Markdown."""
    if not os.path.exists(docx_path):
        raise FileNotFoundError(f"DOCX file not found: {docx_path}")
    try:
        doc = Document(docx_path)
        text = "\n\n".join([p.text for p in doc.paragraphs if p.text.strip()])
        markdown = f"# {os.path.basename(docx_path)}\n\n{text.strip()}"
        return markdown
    except Exception as e:
        raise RuntimeError(f"Error extracting text from DOCX {docx_path}: {e}")

from pptx import Presentation

def convert_pptx_to_markdown(pptx_path: str) -> str:
    """Extracts text from a PPTX file and converts it to Markdown."""
    if not os.path.exists(pptx_path):
        raise FileNotFoundError(f"PPTX file not found: {pptx_path}")
    try:
        presentation = Presentation(pptx_path)
        text = ""
        for slide in presentation.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    text += shape.text_frame.text + "\n\n"
        markdown = f"# {os.path.basename(pptx_path)}\n\n{text.strip()}"
        return markdown
    except Exception as e:
        raise RuntimeError(f"Error extracting text from PPTX {pptx_path}: {e}")
def convert_txt_to_markdown(txt_path: str) -> str:
    """Reads a TXT file and converts it to Markdown."""
    if not os.path.exists(txt_path):
        raise FileNotFoundError(f"TXT file not found: {txt_path}")
    try:
        with open(txt_path, 'r', encoding='utf-8') as f:
            text = f.read()
        markdown = f"# {os.path.basename(txt_path)}\n\n{text.strip()}"
        return markdown
    except Exception as e:
        raise RuntimeError(f"Error reading TXT file {txt_path}: {e}")
from markdownify import markdownify as md

def convert_html_to_markdown(html_content: str, filename: str) -> str:
    """Converts HTML content into Markdown."""
    markdown = md(html_content)
    return f"# {filename}\n\n{markdown.strip()}"


from bs4 import BeautifulSoup
import json
from config import WEB_SEARCH_SUMMARIZATION_PROMPT

def get_links_from_search_results(results: dict) -> list[str]:
    """
    Extracts links from DuckDuckGo Search results.
    Returns a list of URLs.
    """
    links = []
    if results:
        for item in results:
            link = item['href']
            if link:
                links.append(link)
    return links


def synthesize_web_search_response(query: str, df,  num_results=5) -> tuple[str, str | None]:
    res = df.duckduckgo_search(query, max_results=num_results)

    links = get_links_from_search_results(json.loads(res))
    summarized_results = get_summarized_results(links)

    return summarize_web_results_with_ollama(query, summarized_results)


def get_summarized_results(links: list, query: str = 'query') -> tuple[str, str | None]:
    # Scrape and summarize content from each link
    summaries = []
    headers = {"User-Agent": "Mozilla/5.0"}
    reference = ""
    summary = None
    idx = 0

    for link in links:
        try:
            page = requests.get(link, headers=headers, timeout=8)
            page.raise_for_status()
            soup = BeautifulSoup(page.text, "html.parser")
            # Try to get meta description or first two paragraphs
            meta_desc = soup.find("meta", attrs={"name": "description"})
            
            if meta_desc and meta_desc.get("content"):
                summary = meta_desc["content"]
            else:
                paragraphs = soup.find_all("p")
                summary = " ".join([p.get_text(strip=True) for p in paragraphs[:2]])
            if summary:
                summaries.append(f"[{idx+1}] {summary.strip()}")
                reference += f"Source: {link}\n"
        except Exception as e:
            logger.error(f"[{idx+1}] Could not retrieve summary from {link} ({type(e).__name__})", exc_info=True)
    
    if not summaries:
        return "Could not extract useful information from the top web links.", f"Web search performed for: '{query}', but scraping failed."

    answer = "\n\n".join(summaries)
    thinking_content = f"Web search performed using DuckDuckGo Search for: '{query}'. Top {idx} links scraped and summarized.\n" + reference
    return answer, thinking_content


def summarize_web_results_with_ollama(query: str, web_results: list[str]) -> tuple[str, str | None]:
    """
    Summarizes web search results for a query using Ollama and the WEB_SEARCH_SUMMARIZATION_PROMPT.
    Args:
        query (str): The user's question.
        web_results (list[str]): List of web result snippets, each formatted as '[n] snippet (source: url)'.

    Returns:
        tuple[str, str | None]: (summary, thinking_content)
    """
    global llm
    if not llm:
        logger.error("LLM not initialized, cannot summarize web results.")
        return "Error: The AI model is currently unavailable.", None

    # Format web_results as a single string for the prompt
    formatted_results = "\n".join(web_results)

    try:
        # Build the prompt using your template
        final_prompt = WEB_SEARCH_SUMMARIZATION_PROMPT.format(query=query, web_results=formatted_results)
        logger.info("Sending web search summarization prompt to LLM (Ollama)...")
        logger.debug(f"Web Search Summarization Prompt (Start):\n{final_prompt[:200]}...")

        # Call Ollama
        response_object = llm.invoke(final_prompt)
        full_llm_response = getattr(response_object, 'content', str(response_object))

        # Parse the response to separate thinking and answer
        user_answer, thinking_content = parse_llm_response(full_llm_response)

        if not user_answer and thinking_content:
            user_answer = "[AI response consisted only of reasoning. No final answer provided. See thinking process.]"
        elif not user_answer and not thinking_content:
            user_answer = "[AI Response Processing Error: Empty result after parsing]"

        return user_answer.strip(), thinking_content

    except Exception as e:
        logger.error(f"Web search summarization failed: {e}", exc_info=True)
        return f"Error: Web search summarization failed due to {type(e).__name__}.", None