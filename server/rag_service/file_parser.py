import os
import logging
from langchain_community.document_loaders import PyPDFLoader, TextLoader, Docx2txtLoader
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain_core.documents import Document # Import Document
from pptx import Presentation # pip install python-pptx

logger = logging.getLogger(__name__)

def parse_file(file_path: str) -> str:
    """
    Parses various file types and returns their text content.
    Supports .txt, .pdf, .docx, and .pptx.
    """
    if not os.path.exists(file_path):
        logger.error(f"File not found: {file_path}")
        return ""

    file_extension = os.path.splitext(file_path)[1].lower()
    text_content = ""

    try:
        if file_extension == ".txt":
            logger.info(f"Loading .txt file: {file_path}")
            loader = TextLoader(file_path, encoding="utf-8")
            docs = loader.load()
            text_content = "\n".join([doc.page_content for doc in docs])
        elif file_extension == ".pdf":
            logger.info(f"Loading .pdf file: {file_path}")
            loader = PyPDFLoader(file_path)
            docs = loader.load()
            text_content = "\n".join([doc.page_content for doc in docs])
        elif file_extension == ".docx":
            logger.info(f"Loading .docx file: {file_path}")
            loader = Docx2txtLoader(file_path)
            docs = loader.load()
            text_content = "\n".join([doc.page_content for doc in docs])
        elif file_extension == ".pptx":
            logger.info(f"Loading .pptx file: {file_path}")
            prs = Presentation(file_path)
            full_text = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        full_text.append(shape.text)
            text_content = "\n".join(full_text)
        else:
            logger.warning(f"Unsupported file type for parsing: {file_extension} in {file_path}")
            return ""

        if not text_content.strip():
            logger.warning(f"No text content extracted from {file_path}.")
            return ""

        logger.info(f"Successfully extracted {len(text_content)} characters from {os.path.basename(file_path)}.")
        return text_content

    except Exception as e:
        logger.error(f"Error parsing file {file_path}: {e}", exc_info=True)
        return ""

def chunk_text(text_content: str, filename: str, user_id: str) -> list[Document]:
    """
    Chunks raw text content into Langchain Document objects with metadata.
    """
    if not text_content or not text_content.strip():
        logger.warning(f"No text content provided for chunking from file: {filename}. Skipping.")
        return []

    from rag_service import config # Import config here to avoid circular dependency if file_parser imports config earlier

    text_splitter = RecursiveCharacterTextSplitter(
        chunk_size=config.CHUNK_SIZE,
        chunk_overlap=config.CHUNK_OVERLAP,
        length_function=len,
        add_start_index=True,
    )

    chunks = text_splitter.create_documents([text_content])

    # Add metadata to each chunk
    langchain_docs = []
    for i, chunk in enumerate(chunks):
        chunk.metadata = {
            "source": filename,
            "user_id": user_id,
            "chunk_id": f"{user_id}_{filename}_{i}",
            "start_index": chunk.metadata.get("start_index", 0),
        }
        langchain_docs.append(chunk)

    logger.info(f"Chunked '{filename}' into {len(langchain_docs)} documents for user '{user_id}'.")
    return langchain_docs