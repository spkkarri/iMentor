# server/rag_service/file_parser.py
import os
import logging
import re

# --- Library Imports with Graceful Fallbacks ---

try:
    import pypdf
except ImportError:
    print("pypdf not found, PDF parsing will fail. Install with: pip install pypdf")
    pypdf = None

try:
    from docx import Document as DocxDocument
except ImportError:
    print("python-docx not found, DOCX parsing will fail. Install with: pip install python-docx")
    DocxDocument = None

try:
    from pptx import Presentation
    PPTX_SUPPORTED = True
except ImportError:
    print("python-pptx not found, PPTX parsing will fail. Install with: pip install python-pptx")
    PPTX_SUPPORTED = False

# Add this block near the top with the other imports
try:
    from youtube_transcript_api import YouTubeTranscriptApi
    YOUTUBE_SUPPORTED = True
except ImportError:
    print("youtube-transcript-api not found, YouTube parsing will fail. Install with: pip install youtube-transcript-api")
    YouTubeTranscriptApi = None
    YOUTUBE_SUPPORTED = False

from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain_core.documents import Document as LangchainDocument
from . import config # Import from the current package (ai_core_service)

# --- Logger Configuration ---
logger = logging.getLogger(__name__)
logger.setLevel(logging.INFO) # Or DEBUG for more details
handler = logging.StreamHandler()
formatter = logging.Formatter('%(asctime)s - %(name)s - %(levelname)s - %(message)s')
handler.setFormatter(formatter)
if not logger.hasHandlers():
    logger.addHandler(handler)


# --- YouTube Parsing Functions ---

# Add this new function
def _extract_video_id(url):
    """Extracts the YouTube video ID from a URL."""
    if not url: return None
    # Standard watch URL: https://www.youtube.com/watch?v=VIDEO_ID
    match = re.search(r"watch\?v=([^&]+)", url)
    if match:
        return match.group(1)
    # Shortened URL: https://youtu.be/VIDEO_ID
    match = re.search(r"youtu\.be/([^&?/]+)", url)
    if match:
        return match.group(1)
    # Embedded URL: https://www.youtube.com/embed/VIDEO_ID
    match = re.search(r"embed/([^&?/]+)", url)
    if match:
        return match.group(1)
    return None

# Add this new function
def parse_youtube_url(url):
    """Fetches the transcript for a YouTube URL."""
    if not YOUTUBE_SUPPORTED:
        logger.warning("Skipping YouTube URL parsing as youtube-transcript-api is not installed.")
        return None
    
    video_id = _extract_video_id(url)
    if not video_id:
        logger.error(f"Could not extract a valid YouTube video ID from URL: {url}")
        return None
    
    try:
        logger.info(f"Fetching transcript for YouTube video ID: {video_id}")
        transcript_list = YouTubeTranscriptApi.get_transcript(video_id)
        
        # Combine the 'text' fields from the transcript list into a single string
        transcript_text = " ".join([item['text'] for item in transcript_list])
        
        # Replace multiple spaces/newlines with a single space for clean-up
        transcript_text = re.sub(r'\s+', ' ', transcript_text)
        
        logger.debug(f"Extracted {len(transcript_text)} characters from YouTube transcript.")
        return transcript_text.strip() if transcript_text.strip() else None
    except Exception as e:
        logger.error(f"Failed to fetch or process transcript for video ID {video_id}: {e}", exc_info=True)
        return None


# --- File-Based Parsing Functions ---

def parse_pdf(file_path):
    """Extracts text content from a PDF file using pypdf."""
    if not pypdf: return None # Check if library loaded
    text = ""
    try:
        reader = pypdf.PdfReader(file_path)
        num_pages = len(reader.pages)
        for i, page in enumerate(reader.pages):
            try:
                page_text = page.extract_text()
                if page_text:
                    text += page_text + "\n" # Add newline between pages
            except Exception as page_err:
                 logger.warning(f"Error extracting text from page {i+1} of {os.path.basename(file_path)}: {page_err}")
        return text.strip() if text.strip() else None # Return None if empty after stripping
    except FileNotFoundError:
        logger.error(f"PDF file not found: {file_path}")
        return None
    except pypdf.errors.PdfReadError as pdf_err:
        logger.error(f"Error reading PDF {os.path.basename(file_path)} (possibly corrupted or encrypted): {pdf_err}")
        return None
    except Exception as e:
        logger.error(f"Unexpected error parsing PDF {os.path.basename(file_path)}: {e}", exc_info=True)
        return None

def parse_docx(file_path):
    """Extracts text content from a DOCX file."""
    if not DocxDocument: return None # Check if library loaded
    try:
        doc = DocxDocument(file_path)
        text = "\n".join([para.text for para in doc.paragraphs if para.text.strip()])
        return text.strip() if text.strip() else None
    except Exception as e:
        logger.error(f"Error parsing DOCX {os.path.basename(file_path)}: {e}", exc_info=True)
        return None

def parse_txt(file_path):
    """Reads text content from a TXT file (or similar plain text like .py, .js)."""
    try:
        with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
            text = f.read()
        return text.strip() if text.strip() else None
    except Exception as e:
        logger.error(f"Error parsing TXT {os.path.basename(file_path)}: {e}", exc_info=True)
        return None

def parse_pptx(file_path):
    """Extracts text content from a PPTX file."""
    if not PPTX_SUPPORTED:
        logger.warning(f"Skipping PPTX file {os.path.basename(file_path)} as python-pptx is not installed.")
        return None
    text = ""
    try:
        prs = Presentation(file_path)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    shape_text = shape.text.strip()
                    if shape_text:
                        text += shape_text + "\n" # Add newline between shape texts
        return text.strip() if text.strip() else None
    except Exception as e:
        logger.error(f"Error parsing PPTX {os.path.basename(file_path)}: {e}", exc_info=True)
        return None

def parse_file(file_path):
    """Parses a file based on its extension, returning text content or None."""
    _, ext = os.path.splitext(file_path)
    ext = ext.lower()
    logger.debug(f"Attempting to parse file: {os.path.basename(file_path)} (Extension: {ext})")

    if ext == '.pdf':
        return parse_pdf(file_path)
    elif ext == '.docx':
        return parse_docx(file_path)
    elif ext == '.pptx':
        return parse_pptx(file_path)
    elif ext in ['.txt', '.py', '.js', '.md', '.log', '.csv', '.html', '.xml', '.json']:
        return parse_txt(file_path)
    elif ext == '.doc':
        logger.warning(f"Parsing for legacy .doc files is not implemented: {os.path.basename(file_path)}")
        return None
    else:
        logger.warning(f"Unsupported file extension for parsing: {ext} ({os.path.basename(file_path)})")
        return None

# --- Text Processing and Utility Functions ---

def chunk_text(text, file_name, user_id):
    """Chunks text and creates Langchain Documents with metadata."""
    if not text or not isinstance(text, str):
        logger.warning(f"Invalid text input for chunking (file: {file_name}). Skipping.")
        return []

    text_splitter = RecursiveCharacterTextSplitter(
        chunk_size=config.CHUNK_SIZE,
        chunk_overlap=config.CHUNK_OVERLAP,
        length_function=len,
        is_separator_regex=False,
    )

    try:
        chunks = text_splitter.split_text(text)
        if not chunks:
             logger.warning(f"Text splitting resulted in zero chunks for file: {file_name}")
             return []

        documents = []
        for i, chunk in enumerate(chunks):
             if chunk and chunk.strip():
                 documents.append(
                     LangchainDocument(
                         page_content=chunk,
                         metadata={
                             'userId': user_id,
                             'documentName': file_name,
                             'chunkIndex': i
                         }
                     )
                 )
        if documents:
            logger.info(f"Split '{file_name}' into {len(documents)} non-empty chunks.")
        else:
            logger.warning(f"No non-empty chunks created for file: {file_name} after splitting.")
        return documents
    except Exception as e:
        logger.error(f"Error during text splitting for file {file_name}: {e}", exc_info=True)
        return []

def extract_headings(text, chapter_prefix="5."):
    """
    Extracts structured headings from a document like:
    5.1 Air Pollution, 5.2.1 Water Pollution, etc.
    """
    pattern = re.compile(rf"{re.escape(chapter_prefix)}\d*(?:\.\d+)*\s+([A-Za-z].+)")
    matches = pattern.findall(text)
    return list(dict.fromkeys([match.strip() for match in matches]))  # Unique, ordered


# --- High-Level Dispatcher Function ---

# Add this new function at the end of the file
def extract_text_from_input(input_data, input_type, temp_file_path=None):
    """
    A high-level dispatcher to extract text from various input types.
    - input_data: Can be raw text, a URL, or a file object.
    - input_type: 'raw_text', 'youtube_url', or 'file'.
    - temp_file_path: The path where an uploaded file has been saved.
    """
    logger.info(f"Extracting text for input_type: {input_type}")

    if input_type == 'raw_text':
        # Input is already text, just return it after cleaning
        return input_data.strip() if input_data and input_data.strip() else None
    
    elif input_type == 'youtube_url':
        return parse_youtube_url(input_data)
        
    elif input_type == 'file':
        if not temp_file_path:
            logger.error("Input type is 'file' but no temp_file_path was provided.")
            return None
        return parse_file(temp_file_path) # Use your existing file parser
    
    else:
        logger.warning(f"Unsupported input_type specified: {input_type}")
        return None