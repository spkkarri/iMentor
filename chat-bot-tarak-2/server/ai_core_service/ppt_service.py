# FusedChatbot/server/ai_core_service/ppt_service.py

import os
import uuid
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN

# --- Define a Professional LIGHT THEME Color Palette ---
PPT_THEME = {
    "background": RGBColor(0xFF, 0xFF, 0xFF),  # Pure White
    "text": RGBColor(0x21, 0x21, 0x21),        # Very Dark Gray (almost black)
    "accent": RGBColor(0x0D, 0x47, 0xA1)       # Deep, professional Blue
}

_PPT_OUTPUT_DIR = os.path.join(os.path.dirname(__file__), '..', 'generated_ppts')


def apply_master_styles(prs):
    """Applies the color theme to the slide master for consistency."""
    slide_master = prs.slide_masters[0]
    background = slide_master.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = PPT_THEME["background"]

    # Set default text color for all placeholders
    for layout in slide_master.slide_layouts:
        for shape in layout.placeholders:
            if shape.has_text_frame:
                shape.text_frame.paragraphs[0].font.color.rgb = PPT_THEME["text"]


def create_presentation(content: dict) -> str:
    """
    Creates a professionally styled and visually appealing .pptx file with a light theme.
    """
    os.makedirs(_PPT_OUTPUT_DIR, exist_ok=True)

    prs = Presentation()
    prs.slide_width = Inches(16)
    prs.slide_height = Inches(9)

    apply_master_styles(prs)

    # --- Slide 1: Title Slide (Layout 0) ---
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]

    title.text = content.get("title", "AI Generated Presentation")
    p = title.text_frame.paragraphs[0]
    p.font.color.rgb = PPT_THEME["accent"]
    p.font.name = 'Calibri (Headings)'
    p.font.size = Pt(60)
    p.font.bold = True

    subtitle.text = f"{content.get('topic', 'a selected topic')}"
    p = subtitle.text_frame.paragraphs[0]
    p.font.color.rgb = PPT_THEME["text"]
    p.font.name = 'Calibri (Body)'
    p.font.size = Pt(24)

    # --- Content Slides (Layout 1: Title and Content) ---
    slide_layout = prs.slide_layouts[1]
    
    for slide_content in content.get("slides", []):
        slide = prs.slides.add_slide(slide_layout)
        
        title_shape = slide.shapes.title
        body_shape = slide.placeholders[1]

        # Style Slide Title
        title_shape.text = slide_content.get("title", "")
        p = title_shape.text_frame.paragraphs[0]
        p.font.color.rgb = PPT_THEME["accent"]
        p.font.size = Pt(44)
        p.font.bold = True
        
        # --- THIS IS THE CORRECTED BODY SHAPE LOGIC ---
        # Position and size the body placeholder to fill the slide properly
        body_shape.left = Inches(1)
        body_shape.top = Inches(1.75)
        body_shape.width = Inches(14)
        body_shape.height = Inches(5.5)

        tf = body_shape.text_frame
        tf.clear()
        tf.word_wrap = True # Ensure text wraps within the shape

        for bullet_text in slide_content.get("bullets", []):
            p = tf.add_paragraph()
            p.text = bullet_text
            p.font.size = Pt(28)
            p.font.name = 'Calibri'
            p.level = 0
            # Add appropriate line spacing for readability
            p.line_spacing = 1.5
            p.space_after = Pt(10)

    # --- Final Slide (Layout 2: Section Header) ---
    slide_layout = prs.slide_layouts[2]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]

    title.text = "Thank You"
    p = title.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.color.rgb = PPT_THEME["accent"]
    p.font.size = Pt(60)
    p.font.bold = True
    
    subtitle.text = "Questions & Discussion"
    p = subtitle.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.color.rgb = PPT_THEME["text"]
    p.font.size = Pt(32)
    
    # --- Save the presentation ---
    file_id = f"{uuid.uuid4()}.pptx"
    save_path = os.path.join(_PPT_OUTPUT_DIR, file_id)
    
    prs.save(save_path)
    
    return file_id