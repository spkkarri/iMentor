# FusedChatbot/server/scripts/collect_document_data.py

import os
import json
import logging
import docx
import fitz  # PyMuPDF
from pptx import Presentation
from dotenv import load_dotenv
from pymongo import MongoClient

# --- Configuration ---
# Configure logging for clear output
logging.basicConfig(level=logging.INFO, format='%(asctime)s - %(levelname)s - %(message)s')

# Load .env file from the project root to get the MONGO_URI
# This makes the script runnable from different locations (e.g., project root or /server)
try:
    # Try loading from the standard project root location
    project_root = os.path.abspath(os.path.join(os.path.dirname(__file__), '..', '..'))
    dotenv_path = os.path.join(project_root, '.env')
    if not os.path.exists(dotenv_path):
        # Fallback for running from the 'server' directory
        dotenv_path = os.path.join(os.getcwd(), '.env')
    load_dotenv(dotenv_path=dotenv_path)
    logging.info(f"Loaded .env file from: {dotenv_path}")
except Exception as e:
    logging.warning(f"Could not load .env file: {e}. Relying on environment variables.")


MONGO_URI = os.getenv("MONGO_URI")
if not MONGO_URI:
    raise ValueError("MONGO_URI not found in environment variables. Please check your .env file.")

# Define the base path to where user files are stored.
# Path is relative to this script's location.
ASSETS_DIR = os.path.abspath(os.path.join(os.path.dirname(__file__), '..', '..', 'assets'))

# --- File Reading Functions ---

def read_txt_file(file_path: str) -> str:
    """Reads content from a plain text or markdown file."""
    try:
        with open(file_path, 'r', encoding='utf-8') as f:
            return f.read()
    except Exception as e:
        logging.error(f"Error reading TXT/MD file {os.path.basename(file_path)}: {e}")
        return ""

def read_docx_file(file_path: str) -> str:
    """Reads content from a .docx file."""
    try:
        doc = docx.Document(file_path)
        return "\n".join([para.text for para in doc.paragraphs if para.text])
    except Exception as e:
        logging.error(f"Error reading DOCX file {os.path.basename(file_path)}: {e}")
        return ""

def read_pdf_file(file_path: str) -> str:
    """Reads text content from a .pdf file."""
    try:
        with fitz.open(file_path) as doc:
            return "".join([page.get_text() for page in doc])
    except Exception as e:
        logging.error(f"Error reading PDF file {os.path.basename(file_path)}: {e}")
        return ""

def read_pptx_file(file_path: str) -> str:
    """Reads text content from a .pptx file."""
    try:
        prs = Presentation(file_path)
        text_runs = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text_runs.append(shape.text)
        return "\n".join(text_runs)
    except Exception as e:
        logging.error(f"Error reading PPTX file {os.path.basename(file_path)}: {e}")
        return ""

# --- Main Data Collection Logic ---

def collect_document_content() -> list[str]:
    """
    Walks through the ASSETS_DIR, finds all user documents,
    and extracts their text content using appropriate readers.
    """
    logging.info(f"Starting document collection from base directory: {ASSETS_DIR}")
    document_texts = []
    
    if not os.path.exists(ASSETS_DIR):
        logging.warning(f"Assets directory not found at {ASSETS_DIR}. No documents will be collected.")
        return document_texts

    # Map file extensions to their respective reading functions
    readers = {
        ".txt": read_txt_file,
        ".md": read_txt_file,  # Treat markdown as plain text
        ".docx": read_docx_file,
        ".pdf": read_pdf_file,
        ".pptx": read_pptx_file,
    }
    
    # Expecting a structure like: assets/<user>/docs/
    for user_dir in os.listdir(ASSETS_DIR):
        user_docs_path = os.path.join(ASSETS_DIR, user_dir, 'docs')
        
        if os.path.isdir(user_docs_path):
            logging.info(f"Scanning document folder for user '{user_dir}': {user_docs_path}")
            for filename in os.listdir(user_docs_path):
                file_ext = os.path.splitext(filename)[1].lower()
                
                if file_ext in readers:
                    file_path = os.path.join(user_docs_path, filename)
                    logging.info(f"  -> Reading file: {filename}")
                    content = readers[file_ext](file_path)
                    if content:
                        document_texts.append(content)
                    else:
                        logging.warning(f"  -> No content extracted from: {filename}")
    
    logging.info(f"Successfully collected content from {len(document_texts)} documents.")
    return document_texts

def collect_chat_history():
    """
    Connects to MongoDB and collects user-model conversation pairs.
    """
    logging.info("Starting chat history collection from MongoDB...")
    chat_pairs = []
    try:
        client = MongoClient(MONGO_URI)
        # The database name is part of the MONGO_URI, but if not, you'd specify it here.
        # Example: db = client['chatbotGeminiDB6']
        db = client.get_database() # This gets the default DB from the URI
        chathistories = db.chathistories # The collection name is pluralized by mongoose
        
        logging.info(f"Connected to DB: '{db.name}'. Found {chathistories.count_documents({})} chat history documents.")

        for session in chathistories.find():
            messages = session.get('messages', [])
            # We iterate through messages looking for a 'user' message
            # followed immediately by a 'model' message.
            for i in range(len(messages) - 1):
                if messages[i].get('role') == 'user' and messages[i+1].get('role') == 'model':
                    user_message = messages[i].get('parts', [{}])[0].get('text', '').strip()
                    model_message = messages[i+1].get('parts', [{}])[0].get('text', '').strip()
                    
                    # We only want to train on meaningful exchanges
                    if user_message and model_message:
                        chat_pairs.append({
                            "prompt": user_message,
                            "response": model_message
                        })
        
        logging.info(f"Successfully collected {len(chat_pairs)} user-model conversation pairs.")
        return chat_pairs

    except Exception as e:
        logging.error(f"Failed to connect to or read from MongoDB: {e}")
        return []

# --- Script Execution ---

if __name__ == '__main__':
    logging.info("--- Starting Full Data Collection Script ---")
    
    # Collect data from both sources
    documents = collect_document_content()
    chats = collect_chat_history()
    
    # Combine into a single data object
    raw_data = {
        "documents": documents,
        "chat_pairs": chats
    }
    
    # Define the output file path
    output_path = os.path.join(os.path.dirname(__file__), 'raw_data.json')
    
    # Save the raw data to a JSON file
    try:
        with open(output_path, 'w', encoding='utf-8') as f:
            json.dump(raw_data, f, indent=2, ensure_ascii=False)
        logging.info(f"Successfully saved all raw data to {output_path}")
    except Exception as e:
        logging.error(f"Failed to write raw data to JSON file: {e}")

    print("\n--- Final Summary ---")
    print(f"Total documents processed: {len(documents)}")
    print(f"Total chat pairs collected: {len(chats)}")